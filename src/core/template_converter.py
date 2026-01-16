# Copyright 2024 PPT Generate Project
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""PPTX template converter for extracting styles to Marp CSS."""

import os
from pathlib import Path
from typing import Optional, Dict, Any, Tuple
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

from src.utils.logging import Logger


logger = Logger.get_logger(__name__)


class PPTXTemplateConverter:
    """Convert PPTX templates to Marp CSS styles."""

    def __init__(self) -> None:
        """Initialize the converter."""
        pass

    def convert_pptx_to_marp_css(self, pptx_path: str) -> str:
        """
        Convert PPTX file to Marp CSS.

        Args:
            pptx_path: Path to the PPTX file

        Returns:
            CSS string compatible with Marp
        """
        try:
            logger.info(f"Converting PPTX template: {pptx_path}")

            # Load presentation
            prs = Presentation(pptx_path)

            # Extract theme information
            theme_info = self._extract_theme_info(prs)

            # Generate CSS
            css = self._generate_marp_css(theme_info)

            logger.info("PPTX to Marp CSS conversion completed")
            return css

        except Exception as e:
            logger.error(f"PPTX conversion failed: {str(e)}")
            raise

    def _extract_theme_info(self, presentation: Presentation) -> Dict[str, Any]:
        """
        Extract theme information from PPTX presentation.

        Args:
            presentation: Loaded PPTX presentation

        Returns:
            Dictionary containing extracted theme information
        """
        theme_info = {
            "colors": {},
            "fonts": {},
            "backgrounds": [],
            "layouts": []
        }

        try:
            # Extract slide master information
            if presentation.slide_master:
                master = presentation.slide_master

                # Extract theme colors
                theme_info["colors"] = self._extract_theme_colors(master)

                # Extract fonts
                theme_info["fonts"] = self._extract_theme_fonts(master)

                # Extract background information
                theme_info["backgrounds"] = self._extract_background_info(master)

            # Extract layout information from first few slides
            theme_info["layouts"] = self._analyze_layouts(presentation)

        except Exception as e:
            logger.warning(f"Error extracting theme info: {str(e)}")

        return theme_info

    def _extract_theme_colors(self, slide_master) -> Dict[str, str]:
        """Extract theme colors from slide master."""
        colors = {}

        try:
            # Get theme color scheme
            theme = slide_master.theme

            # Extract all main theme colors
            color_mapping = {
                "primary": MSO_THEME_COLOR.ACCENT_1,
                "secondary": MSO_THEME_COLOR.ACCENT_2,
                "tertiary": MSO_THEME_COLOR.ACCENT_3,
                "accent4": MSO_THEME_COLOR.ACCENT_4,
                "accent5": MSO_THEME_COLOR.ACCENT_5,
                "accent6": MSO_THEME_COLOR.ACCENT_6,
                "text": MSO_THEME_COLOR.TEXT_1,
                "text2": MSO_THEME_COLOR.TEXT_2,
                "background": MSO_THEME_COLOR.BACKGROUND_1,
                "background2": MSO_THEME_COLOR.BACKGROUND_2,
                "hyperlink": MSO_THEME_COLOR.HYPERLINK,
                "followed_hyperlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK
            }

            for color_name, theme_color in color_mapping.items():
                try:
                    # Try to get the mapped color from theme
                    color = slide_master.theme.theme_color_scheme.rgb_color(theme_color)
                    colors[color_name] = self._rgb_to_hex(color)
                except Exception as e:
                    logger.debug(f"Could not extract theme color {color_name}: {str(e)}")
                    # Standard fallbacks for missing colors
                    if "text" in color_name: colors[color_name] = "#000000"
                    elif "background" in color_name: colors[color_name] = "#FFFFFF"
                    else: colors[color_name] = "#2E75B6"

        except Exception as e:
            logger.warning(f"Error extracting theme colors: {str(e)}")
            # Minimal defaults
            colors = {
                "primary": "#2E75B6",
                "text": "#000000",
                "background": "#FFFFFF"
            }

        return colors

    def _extract_theme_fonts(self, slide_master) -> Dict[str, str]:
        """Extract theme fonts from slide master."""
        fonts = {}

        try:
            theme = slide_master.theme

            # Extract major fonts
            try:
                fonts["heading"] = theme.font_scheme.major_font.latin.typeface
            except:
                fonts["heading"] = "Arial"

            try:
                fonts["body"] = theme.font_scheme.minor_font.latin.typeface
            except:
                fonts["body"] = "Arial"

        except Exception as e:
            logger.warning(f"Error extracting theme fonts: {str(e)}")
            fonts = {
                "heading": "Arial",
                "body": "Arial"
            }

        return fonts

    def _extract_background_info(self, slide_master) -> list:
        """Extract background information from slide master."""
        backgrounds = []

        try:
            # Check if slide master has background
            if hasattr(slide_master, 'background'):
                bg = slide_master.background
                if bg.fill.type == 1:  # Solid fill
                    try:
                        rgb_color = bg.fill.fore_color.rgb_color
                        backgrounds.append({
                            "type": "solid",
                            "color": self._rgb_to_hex(rgb_color)
                        })
                    except:
                        pass

        except Exception as e:
            logger.warning(f"Error extracting background info: {str(e)}")

        return backgrounds

    def _analyze_layouts(self, presentation: Presentation) -> list:
        """Analyze common layouts used in the presentation."""
        layouts = []

        try:
            # Analyze first few slides to understand layout patterns
            for i, slide in enumerate(presentation.slides[:5]):
                layout_info = {
                    "title_placeholders": 0,
                    "content_placeholders": 0,
                    "has_images": False
                }

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if shape.text.startswith("Title") or len(shape.text) < 50:
                            layout_info["title_placeholders"] += 1
                        else:
                            layout_info["content_placeholders"] += 1

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        layout_info["has_images"] = True

                layouts.append(layout_info)

        except Exception as e:
            logger.warning(f"Error analyzing layouts: {str(e)}")

        return layouts

    def _rgb_to_hex(self, rgb_color: RGBColor) -> str:
        """Convert RGB color to hex string."""
        return f"#{rgb_color.rgb[0]:02x}{rgb_color.rgb[1]:02x}{rgb_color.rgb[2]:02x}"

    def _generate_marp_css(self, theme_info: Dict[str, Any]) -> str:
        """
        Generate Marp-compatible CSS from extracted theme information.

        Args:
            theme_info: Dictionary containing theme information

        Returns:
            CSS string for Marp
        """
        colors = theme_info.get("colors", {})
        fonts = theme_info.get("fonts", {})

        css = f"""/* Generated Marp CSS from PPTX template */

:root {{
  --color-primary: {colors.get('primary', '#2E75B6')};
  --color-secondary: {colors.get('secondary', '#5B9BD5')};
  --color-tertiary: {colors.get('tertiary', '#A8D08D')};
  --color-accent4: {colors.get('accent4', '#ED7D31')};
  --color-accent5: {colors.get('accent5', '#4472C4')};
  --color-accent6: {colors.get('accent6', '#70AD47')};
  --color-text: {colors.get('text', '#000000')};
  --color-text-dim: {colors.get('text2', '#666666')};
  --color-background: {colors.get('background', '#FFFFFF')};
  --color-background-alt: {colors.get('background2', '#F2F2F2')};
  --color-highlight: {colors.get('hyperlink', '#0563C1')};

  --font-heading: '{fonts.get('heading', 'Arial')}', sans-serif;
  --font-body: '{fonts.get('body', 'Arial')}', sans-serif;
}}

/* Base slide styles */
section {{
  background-color: var(--color-background);
  color: var(--color-text);
  font-family: var(--font-body);
  display: flex;
  flex-direction: column;
  padding: 40px;
  background-image: linear-gradient(to bottom right, var(--color-background), var(--color-background-alt));
}}

/* Heading styles */
h1, h2, h3, h4, h5, h6 {{
  color: var(--color-primary);
  font-family: var(--font-heading);
  font-weight: bold;
  margin-top: 0;
}}

h1 {{
  font-size: 2.8em;
  margin-bottom: 0.6em;
  border-bottom: 3px solid var(--color-primary);
  padding-bottom: 0.2em;
}}

h2 {{
  font-size: 2.2em;
  margin-bottom: 0.5em;
  color: var(--color-secondary);
}}

h3 {{
  font-size: 1.8em;
  margin-bottom: 0.4em;
  color: var(--color-tertiary);
}}

/* List styles */
ul, ol {{
  color: var(--color-text);
  font-size: 1.2em;
  line-height: 1.6;
}}

li {{
  margin-bottom: 0.6em;
}}

li::marker {{
  color: var(--color-primary);
}}

/* Link styles */
a {{
  color: var(--color-highlight);
  text-decoration: underline;
}}

/* Code styles */
code {{
  background-color: var(--color-background-alt);
  color: var(--color-accent4);
  padding: 0.2em 0.4em;
  border-radius: 6px;
  font-family: 'Consolas', 'Monaco', monospace;
}}

pre code {{
  background-color: transparent;
  padding: 0;
}}

/* Blockquote styles */
blockquote {{
  border-left: 8px solid var(--color-primary);
  padding: 10px 20px;
  margin: 20px 0;
  background-color: var(--color-background-alt);
  color: var(--color-text-dim);
  font-style: italic;
  border-radius: 0 10px 10px 0;
}}

/* Custom classes for different layouts */
.split {{
  display: flex;
  flex-direction: row;
  align-items: center;
  gap: 40px;
  flex: 1;
}}

.split > * {{
  flex: 1;
}}

.cover {{
  display: flex;
  flex-direction: column;
  justify-content: center;
  align-items: center;
  text-align: center;
  background-color: var(--color-primary);
  color: var(--color-background);
}}

.cover h1 {{
  font-size: 4.5em;
  margin-bottom: 0.2em;
  color: var(--color-background);
  border-bottom: none;
}}

.cover p {{
  font-size: 1.8em;
  opacity: 0.9;
}}

.quote {{
  display: flex;
  flex-direction: column;
  justify-content: center;
  text-align: center;
  font-style: italic;
  font-size: 2.2em;
  padding: 60px;
  color: var(--color-secondary);
}}

.quote::before {{
  content: '“';
  font-size: 3em;
  line-height: 0.1em;
  margin-right: 0.25em;
  vertical-align: -0.4em;
  color: var(--color-tertiary);
}}

footer {{
  font-size: 0.6em;
  color: var(--color-text-dim);
  position: absolute;
  bottom: 20px;
  right: 40px;
}}
"""

        return css