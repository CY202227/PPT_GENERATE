#!/usr/bin/env python
"""
PowerPoint生成器
支持基于模板和大纲生成新的PPT内容
"""
import os
import json
import logging
import uuid
import shutil
import zipfile
import xml.etree.ElementTree as ET

from openai import OpenAI
from typing import Dict, Any, Tuple, Optional, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.shapes.group import GroupShape
import traceback

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)



class PPTGenerator:
    
    def __init__(self, api_key: str, base_url: str):
        try:
            self.client = OpenAI(api_key=api_key, base_url=base_url)
            logger.info("OpenAI客户端初始化成功")
        except Exception as e:
            logger.error(f"初始化OpenAI客户端失败: {e}")
            raise
    
    def detect_slide_type(self, slide_xml_path: str) -> Tuple[bool, bool]:
        """
        检测幻灯片是否包含SmartArt或图表
        
        Args:
            slide_xml_path: 幻灯片XML文件的路径
            
        Returns:
            Tuple[bool, bool]: (has_smartart, has_chart)
        """
        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
                'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
            }
            
            # 检查SmartArt
            # SmartArt在PPT XML中通常由 <dgm:relIds> 标签标识
            has_smartart = len(root.findall('.//dgm:relIds', namespaces)) > 0
            
            # 检查图表
            # 图表在PPT XML中通常由 <c:chart> 标签标识
            has_chart = len(root.findall('.//c:chart', namespaces)) > 0
            
            logger.info(f"幻灯片类型检测结果 - SmartArt: {has_smartart}, Chart: {has_chart}")
            return has_smartart, has_chart
            
        except ET.ParseError as e:
            logger.error(f"解析XML文件失败: {e}")
            return False, False
        except Exception as e:
            logger.error(f"检测幻灯片类型时发生错误: {e}")
            return False, False
    
    def extract_pptx(self, pptx_path: str, extract_dir: str) -> bool:
        """解压PPTX文件"""
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            logger.info(f"成功解压PPTX到: {extract_dir}")
            return True
        except Exception as e:
            logger.error(f"解压PPTX失败: {e}")
            return False

    def create_pptx(self, source_dir: str, output_path: str) -> bool:
        """将文件夹重新打包为PPTX"""
        try:
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for root, _, files in os.walk(source_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, source_dir)
                        zipf.write(file_path, arcname)
            logger.info(f"成功创建PPTX: {output_path}")
            return True
        except Exception as e:
            logger.error(f"创建PPTX失败: {e}")
            return False
    
    def update_chart_text(self, slide_xml_path: str, text_content: str) -> bool:
        """
        更新图表中的文字内容
        
        Args:
            slide_xml_path: 幻灯片XML文件的路径
            text_content: 新的文本内容
            
        Returns:
            bool: 是否成功更新
        """
        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }
            
            # 查找所有图表引用
            chart_refs = root.findall('.//c:chart', namespaces)
            if not chart_refs:
                return False
                
            # 获取幻灯片所在目录
            slide_dir = os.path.dirname(slide_xml_path)
            ppt_dir = os.path.dirname(slide_dir)
            
            for chart_ref in chart_refs:
                # 获取图表关系ID
                rid = chart_ref.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if not rid:
                    continue
                    
                # 读取关系文件找到图表XML
                rels_path = os.path.join(slide_dir, '_rels', os.path.basename(slide_xml_path) + '.rels')
                if not os.path.exists(rels_path):
                    continue
                    
                rels_tree = ET.parse(rels_path)
                rels_root = rels_tree.getroot()
                
                # 查找图表文件路径
                chart_rel = rels_root.find(f".//Relationship[@Id='{rid}']")
                if chart_rel is None:
                    continue
                    
                chart_path = chart_rel.get('Target')
                if not chart_path:
                    continue
                    
                # 转换为绝对路径
                if not os.path.isabs(chart_path):
                    chart_path = os.path.join(ppt_dir, chart_path.lstrip('/'))
                
                # 更新图表文本
                if os.path.exists(chart_path):
                    chart_tree = ET.parse(chart_path)
                    chart_root = chart_tree.getroot()
                    
                    # 获取图表中的所有文本
                    text_elements = []
                    node_texts = []
                    
                    # 标题
                    title = chart_root.find('.//c:title//c:v', namespaces)
                    if title is not None and title.text:
                        text_elements.append(title)
                        node_texts.append(title.text.strip())
                    
                    # 图例
                    legends = chart_root.findall('.//c:legend//c:v', namespaces)
                    for legend in legends:
                        if legend.text:
                            text_elements.append(legend)
                            node_texts.append(legend.text.strip())
                    
                    # 轴标题和标签
                    axes = chart_root.findall('.//c:axis//c:v', namespaces)
                    for axis in axes:
                        if axis.text:
                            text_elements.append(axis)
                            node_texts.append(axis.text.strip())
                    
                    # 数据标签
                    data_labels = chart_root.findall('.//c:dLbls//c:v', namespaces)
                    for label in data_labels:
                        if label.text:
                            text_elements.append(label)
                            node_texts.append(label.text.strip())
                    
                    if text_elements:
                        # 构建提示词
                        prompt = f"""
你是PPT图表内容生成专家。请根据以下内容生成新的图表文本。

原有文本:
{node_texts}

新的内容:
{text_content}

要求:
1. 必须生成与原文本数量完全相同的文本
2. 每段文本长度要接近原文本
3. 保持原有文本的功能（如果是标题就生成标题，如果是图例就生成图例等）
4. 内容要符合图表展示风格

请返回JSON格式:
{{
  "texts": {node_texts}  // 数组长度必须是{len(node_texts)}
}}
"""
                        # 调用API生成新内容
                        response = self.client.chat.completions.create(
                            model="Qwen-72B",
                            messages=[{"role": "user", "content": prompt}],
                            response_format={"type": "json_object"},
                            temperature=0.7,
                            max_tokens=1000
                        )
                        
                        result = json.loads(response.choices[0].message.content)
                        new_texts = result.get("texts", [])
                        
                        # 确保生成的文本数量正确
                        if len(new_texts) != len(text_elements):
                            logger.warning(f"生成的文本数量({len(new_texts)})与原文本数量({len(text_elements)})不匹配")
                            continue
                        
                        # 更新图表中的文本
                        for elem, new_text in zip(text_elements, new_texts):
                            old_text = elem.text
                            elem.text = new_text
                            logger.info(f"更新图表文本: '{old_text[:30]}...' -> '{new_text[:30]}...'")
                        
                        # 保存修改后的图表XML
                        chart_tree.write(chart_path, encoding='UTF-8', xml_declaration=True)
                        logger.info(f"成功更新图表内容: {chart_path}")
            
            return True
            
        except Exception as e:
            logger.error(f"更新图表文本失败: {e}")
            logger.error(traceback.format_exc())
            return False

    def load_outline(self, outline_path: str) -> Dict[str, Any]:
        """加载大纲文件"""
        try:
            with open(outline_path, 'r', encoding='utf-8') as f:
                outline = json.load(f)
            logger.info("成功加载大纲文件")
            return outline
        except Exception as e:
            logger.error(f"加载大纲文件失败: {e}")
            return {}

    def generate_content_for_slide_type(self, slide_xml_path: str, text_content: str, outline: Dict[str, Any] = None, slide_index: int = 0) -> Dict[str, Any]:
        """根据页面内容生成相应的新内容"""
        try:
            # 首先检测幻灯片类型
            has_smartart, has_chart = self.detect_slide_type(slide_xml_path)
            
            # 如果是图表，更新图表文本
            if has_chart:
                logger.info("检测到图表，开始更新图表文本")
                self.update_chart_text(slide_xml_path, text_content)
            
            # 如果是SmartArt，则跳过内容生成
            if has_smartart:
                logger.info("检测到SmartArt，跳过内容生成")
                return {"texts": []}
            
            # 解析XML文件
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }
            
            # 分析页面结构
            text_boxes = []
            for shape in root.findall('.//p:sp', namespaces):
                # 获取文本框类型（标题、正文等）
                nvSpPr = shape.find('.//p:nvSpPr', namespaces)
                shape_type = "body"
                if nvSpPr is not None:
                    ph = nvSpPr.find('.//p:ph', namespaces)
                    if ph is not None and 'type' in ph.attrib:
                        shape_type = ph.attrib['type']
                
                # 获取位置和大小
                xfrm = shape.find('.//a:xfrm', namespaces)
                position = {"x": 0, "y": 0}
                size = {"width": 0, "height": 0}
                if xfrm is not None:
                    off = xfrm.find('a:off', namespaces)
                    ext = xfrm.find('a:ext', namespaces)
                    if off is not None:
                        position = {
                            "x": int(off.get('x', 0)),
                            "y": int(off.get('y', 0))
                        }
                    if ext is not None:
                        size = {
                            "width": int(ext.get('cx', 0)),
                            "height": int(ext.get('cy', 0))
                        }
                
                # 获取文本内容和样式
                txBody = shape.find('.//p:txBody', namespaces)
                if txBody is not None:
                    text_content = ""
                    is_vertical = False
                    font_size = 0
                    font_name = ""
                    font_color = None
                    
                    # 检查是否是竖排文本
                    bodyPr = txBody.find('a:bodyPr', namespaces)
                    if bodyPr is not None and 'vert' in bodyPr.attrib:
                        is_vertical = bodyPr.attrib['vert'] == 'eaVert'
                    
                    # 获取文本样式
                    for p in txBody.findall('.//a:p', namespaces):
                        for r in p.findall('.//a:r', namespaces):
                            rPr = r.find('a:rPr', namespaces)
                            if rPr is not None:
                                if 'sz' in rPr.attrib:
                                    font_size = max(font_size, int(rPr.get('sz', 0)))
                                if 'typeface' in rPr.attrib:
                                    font_name = rPr.get('typeface', '')
                                solidFill = rPr.find('a:solidFill', namespaces)
                                if solidFill is not None:
                                    srgbClr = solidFill.find('a:srgbClr', namespaces)
                                    if srgbClr is not None and 'val' in srgbClr.attrib:
                                        font_color = srgbClr.get('val')
                    
                    text_boxes.append({
                        "type": shape_type,
                        "position": position,
                        "size": size,
                        "is_vertical": is_vertical,
                        "font_size": font_size,
                        "font_name": font_name,
                        "font_color": font_color,
                        "node": txBody
                    })
            
            # 构建提示词
            outline_info = ""
            if outline:
                current_section = None
                # 根据幻灯片索引确定当前在大纲中的位置
                section_length = len(outline.get("sections", []))
                if section_length > 0:
                    section_index = min(slide_index // 5, section_length - 1)  # 假设每5页对应一个大纲部分
                    current_section = outline["sections"][section_index]
                
                if current_section:
                    outline_info = f"""
当前大纲位置：第{section_index + 1}部分
当前部分主题：{current_section.get('title', '')}
当前部分内容：{json.dumps(current_section, ensure_ascii=False, indent=2)}
"""

            # 构建文本框结构信息
            text_boxes_info = json.dumps([{
                "type": box["type"],
                "is_vertical": box["is_vertical"],
                "font_size": box["font_size"],
                "position": "标题" if box["type"] == "title" else "正文",
                "size": box["size"],
                "style": {
                    "font_size": box["font_size"] // 100 if box["font_size"] else 18,
                    "font_name": "微软雅黑",
                    "font_color": "000000",
                    "text_effects": ["shadow_soft"] if box["type"] == "title" else [],
                    "alignment": "center" if box["type"] == "title" else "left",
                    "line_spacing": 1.2
                }
            } for box in text_boxes], ensure_ascii=False, indent=2)

            # 构建返回格式示例
            return_format_example = json.dumps({
                "texts": [
                    {
                        "content": "新的文本内容",
                        "type": "标题/正文",
                        "style": {
                            "font_size": 24,
                            "font_name": "微软雅黑",
                            "font_color": "000000",
                            "text_effects": ["shadow_soft", "glow_subtle"],
                            "alignment": "center",
                            "line_spacing": 1.2
                        }
                    }
                ]
            }, ensure_ascii=False, indent=2)

            prompt = f"""
你是PPT内容生成专家。请根据页面结构和主题生成新的内容。

页面布局：
{text_boxes_info}

新的内容主题：
{text_content}

{outline_info}

要求：
1. 请为每个文本框生成新的内容
2. 考虑文本框的类型（标题/正文）和样式（竖排/横排）
3. 生成的内容长度要适合文本框大小
4. 内容要符合PPT展示风格
5. 保持整页内容的连贯性和逻辑性
6. 如果是竖排文本，生成的内容要适合竖排展示
7. 完全忽略原有内容，只根据主题和布局生成新内容
8. 标题要简洁有力，正文要详细专业
9. 根据文本框大小和位置生成合适的内容长度
10. 考虑字体大小和颜色的专业搭配

请按照以下格式返回JSON（数组长度必须是{len(text_boxes)}）：
{return_format_example}
"""
            # 调用API生成新内容
            response = self.client.chat.completions.create(
                model="Qwen-72B",
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"},
                temperature=0.7,
                max_tokens=1000
            )
            
            result = json.loads(response.choices[0].message.content)
            new_texts = result.get("texts", [])
            
            # 更新文本内容
            if len(new_texts) == len(text_boxes):
                for box, new_text in zip(text_boxes, new_texts):
                    # 找到所有文本节点
                    for p in box["node"].findall('.//a:p', namespaces):
                        # 设置段落属性
                        pPr = p.find('a:pPr', namespaces)
                        if pPr is None:
                            pPr = ET.SubElement(p, 'a:pPr')
                        
                        # 设置行间距
                        if "line_spacing" in new_text.get("style", {}):
                            spacing = int(float(new_text["style"]["line_spacing"]) * 100000)
                            pPr.set('spcAft', str(spacing))
                            pPr.set('spcBef', str(spacing))
                        
                        # 设置对齐方式
                        if "alignment" in new_text.get("style", {}):
                            algn = new_text["style"]["alignment"]
                            if algn == "center":
                                pPr.set('algn', 'ctr')
                            elif algn == "right":
                                pPr.set('algn', 'r')
                            else:
                                pPr.set('algn', 'l')
                        
                        # 处理每个文本运行
                        for r in p.findall('.//a:r', namespaces):
                            t = r.find('a:t', namespaces)
                            if t is not None:
                                t.text = new_text["content"]
                                
                                # 应用样式
                                style = new_text.get("style", {})
                                rPr = r.find('a:rPr', namespaces)
                                if rPr is not None:
                                    # 设置字体大小
                                    if "font_size" in style:
                                        rPr.set('sz', str(int(style["font_size"]) * 100))
                                    
                                    # 设置字体
                                    if "font_name" in style:
                                        rPr.set('typeface', style["font_name"])
                                    
                                    # 设置颜色
                                    if "font_color" in style:
                                        solidFill = ET.SubElement(rPr, 'a:solidFill')
                                        srgbClr = ET.SubElement(solidFill, 'a:srgbClr')
                                        srgbClr.set('val', style["font_color"])
                                    
                                    # 应用文本效果
                                    if "text_effects" in style:
                                        for effect in style["text_effects"]:
                                            if effect == "shadow_soft":
                                                shadow = ET.SubElement(rPr, 'a:shadow')
                                                shadow.set('blurRad', '50800')
                                                shadow.set('dist', '38100')
                                                shadow.set('dir', '2700000')
                                                shadow.set('algn', 'tl')
                                                shadow.set('rotWithShape', '0')
                                            elif effect == "glow_subtle":
                                                glow = ET.SubElement(rPr, 'a:glow')
                                                glow.set('rad', '38100')
                                                srgbClr = ET.SubElement(glow, 'a:srgbClr')
                                                srgbClr.set('val', '000000')
                                                alpha = ET.SubElement(srgbClr, 'a:alpha')
                                                alpha.set('val', '40000')
                                
                                logger.info(f"更新文本为: '{new_text['content'][:30]}...'")
            
            # 保存修改后的XML
            tree.write(slide_xml_path, encoding='UTF-8', xml_declaration=True)
            logger.info(f"成功更新幻灯片内容")
            return result
            
            logger.warning(f"生成的文本数量({len(new_texts)})与文本框数量({len(text_boxes)})不匹配")
            return {"texts": []}
            
        except Exception as e:
            logger.error(f"生成内容失败: {e}")
            logger.error(traceback.format_exc())
            return {"texts": []}
    
    def smart_template_generation(self, template_path: str, output_path: str, text_content: str, outline_path: str = None) -> bool:
        """智能模板生成：复制模板并修改内容"""
        try:
            # 1. 加载大纲（如果提供）
            outline = None
            if outline_path:
                outline = self.load_outline(outline_path)
                if not outline:
                    return False
            
            # 2. 创建临时目录
            temp_dir = f"temp_{uuid.uuid4().hex}"
            os.makedirs(temp_dir, exist_ok=True)
            
            # 3. 复制模板文件
            temp_pptx = os.path.join(temp_dir, "temp.pptx")
            shutil.copy2(template_path, temp_pptx)
            
            # 4. 解压PPTX
            extract_dir = os.path.join(temp_dir, "extracted")
            if not self.extract_pptx(temp_pptx, extract_dir):
                return False
            
            # 5. 修改幻灯片内容
            slides_dir = os.path.join(extract_dir, "ppt", "slides")
            if os.path.exists(slides_dir):
                # 遍历所有幻灯片XML文件
                for slide_index, slide_file in enumerate(sorted(os.listdir(slides_dir))):
                    if slide_file.endswith(".xml"):
                        slide_path = os.path.join(slides_dir, slide_file)
                        # 为每个幻灯片生成新内容
                        self.generate_content_for_slide_type(slide_path, text_content, outline, slide_index)
            
            # 6. 重新打包为PPTX
            if not self.create_pptx(extract_dir, output_path):
                return False
            
            # 7. 清理临时文件
            shutil.rmtree(temp_dir)
            logger.info("清理临时文件完成")
            
            logger.info(f"PPT生成完成，保存到: {os.path.abspath(output_path)}")
            return True
            
        except Exception as e:
            logger.error(f"PPT生成失败: {e}")
            logger.error(traceback.format_exc())
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
            return False
    
    def test_api_connection(self) -> bool:
        """测试API连接"""
        try:
            logger.info("测试API连接...")
            response = self.client.chat.completions.create(
                model="Qwen-72B",
                messages=[{"role": "user", "content": "你好"}],
                max_tokens=10
            )
            logger.info("API连接测试成功")
            return True
        except Exception as e:
            logger.error(f"API连接测试失败: {e}")
            return False

    def _extract_shape_content(self, shape) -> Dict[str, Any]:
        """递归提取形状内容"""
        element = None
        if isinstance(shape, GroupShape):
            elements = []
            for s in shape.shapes:
                child = self._extract_shape_content(s)
                if child:
                    elements.append(child)
            if not elements:
                return None
            element = { "type": "group", "elements": elements }

        elif hasattr(shape, 'text_frame') and shape.has_text_frame:
            text_frame = shape.text_frame
            content = text_frame.text
            element = { "type": "text", "content": content }

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element = {"type": "image"}

        elif hasattr(shape, 'table'):
            rows = []
            for row in shape.table.rows:
                cell_contents = [cell.text for cell in row.cells]
                rows.append(cell_contents)
            element = {"type": "table", "rows": rows}
        
        if element:
            element["id"] = shape.shape_id
            element["name"] = shape.name
            element["position"] = {"left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height}

        return element

    def extract_slide_content_for_llm(self, slide) -> Dict[str, Any]:
        """将幻灯片转换为结构化的JSON格式给LLM"""
        slide_content = {
            "layout_type": slide.slide_layout.name,
            "elements": []
        }
        for shape in slide.shapes:
            element_data = self._extract_shape_content(shape)
            if element_data:
                slide_content["elements"].append(element_data)
        return slide_content

    def is_ending_slide(self, slide_content: Dict[str, Any]) -> Tuple[bool, Optional[str]]:
        """判断是否是结束页，并返回结束语"""
        ending_keywords = ["感谢", "谢谢", "thanks", "thank you", "the end"]
        
        for element in slide_content.get("elements", []):
            if element.get("type") == "text":
                content = element.get("content", "").lower().strip()
                if len(content) < 15 and any(keyword in content.lower() for keyword in ending_keywords):
                    return True, element.get("content")
                if all(char in "！!。.，,、?？" or any(keyword in content.lower() for keyword in ending_keywords) for char in content):
                    return True, element.get("content")
        return False, None

    def get_section_for_slide_index(self, slide_index: int, total_slides: int, outline: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        """根据幻灯片索引获取对应的章节信息"""
        if slide_index < 2:
            return None
            
        section_keys = [key for key in outline.keys() if key.startswith('section')]
        if not section_keys:
            return None
            
        content_slides = total_slides - 2
        slides_per_section = content_slides / len(section_keys)
        adjusted_index = slide_index - 2
        section_index = int(adjusted_index / slides_per_section)
        section_index = min(section_index, len(section_keys) - 1)
        
        section_key = section_keys[section_index]
        return outline.get(section_key)

    def _get_new_content_map(self, elements: List[Dict[str, Any]]) -> Dict[int, str]:
        """递归遍历生成的elements，创建一个从shape_id到新内容的映射"""
        content_map = {}
        for element in elements:
            shape_id = element.get("id")
            if not shape_id:
                continue
            
            if element.get("type") == "text" and "content" in element:
                content_map[shape_id] = element["content"]
            elif element.get("type") == "table" and "rows" in element:
                content_map[shape_id] = element["rows"]
            
            if element.get("type") == "group" and "elements" in element:
                content_map.update(self._get_new_content_map(element["elements"]))
                
        return content_map

    def _apply_content_to_shapes(self, shapes: List[Any], content_map: Dict[int, Any]):
        """递归遍历幻灯片中的所有形状，并使用content_map更新文本"""
        for shape in shapes:
            if shape.shape_id in content_map:
                new_content = content_map[shape.shape_id]
                
                is_text_shape = hasattr(shape, 'has_text_frame') and shape.has_text_frame
                is_table_shape = hasattr(shape, 'has_table') and shape.has_table

                if is_text_shape and isinstance(new_content, str):
                    text_frame = shape.text_frame
                    if text_frame.paragraphs:
                        p = text_frame.paragraphs[0]
                        if p.runs:
                            p.runs[0].text = new_content
                            for i in range(len(p.runs) - 1, 0, -1):
                                p._p.remove(p.runs[i]._r)
                        else:
                            p.add_run().text = new_content
                        
                        for i in range(len(text_frame.paragraphs) - 1, 0, -1):
                            p_to_remove = text_frame.paragraphs[i]
                            text_frame._txBody.remove(p_to_remove._p)
                
                elif is_table_shape and isinstance(new_content, list):
                    table = shape.table
                    for r, row_content in enumerate(new_content):
                        for c, cell_content in enumerate(row_content):
                            if r < len(table.rows) and c < len(table.columns):
                                cell = table.cell(r, c)
                                if isinstance(cell_content, dict):
                                    cell.text = str(cell_content.get("content", ""))
                                else:
                                    cell.text = str(cell_content)

            if isinstance(shape, GroupShape):
                self._apply_content_to_shapes(shape.shapes, content_map)

    def apply_generated_content(self, slide, generated_content: Dict[str, Any]):
        """将生成的内容应用到幻灯片"""
        content_map = self._get_new_content_map(generated_content.get("elements", []))
        self._apply_content_to_shapes(slide.shapes, content_map)

    def analyze_and_generate_ppt(
        self,
        template_path: str,
        topic: str,
        output_path: str,
        model_name: str,
        outline: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """分析PPT模板并生成新的PPT内容"""
        try:
            if not os.path.exists(template_path):
                return {"error": f"模板文件不存在：{template_path}"}
                
            prs = Presentation(template_path)
            new_prs = Presentation(template_path)
            
            total_slides = len(prs.slides)
            
            # 遍历每个幻灯片
            for slide_index, slide in enumerate(prs.slides):
                logger.info(f"处理第 {slide_index + 1} 张幻灯片")
                
                # 提取幻灯片内容为JSON格式
                slide_content_for_llm = self.extract_slide_content_for_llm(slide)
                
                # 检查是否是最后一页，且是结束语
                is_ending, ending_text = self.is_ending_slide(slide_content_for_llm)
                if slide_index == total_slides - 1 and is_ending and ending_text:
                    logger.info(f"检测到结束页，保持原有结束语：{ending_text}")
                    continue
                
                # 获取当前幻灯片对应的大纲部分
                current_section = None
                outline_info = ""
                if outline:
                    current_section = self.get_section_for_slide_index(slide_index, total_slides, outline)
                    if current_section:
                        outline_info = f"""
当前大纲信息：
- 总标题：{outline.get('title', '')}
- 当前章节：{current_section.get('title', '')}
- 章节内容：{current_section.get('content', '')}
- 日期：{outline.get('date', '')}
"""
                    elif slide_index == 0:
                        outline_info = f"""
封面信息：
- 总标题：{outline.get('title', '')}
- 日期：{outline.get('date', '')}
"""
                    elif slide_index == 1:
                        outline_info = f"""
目录页信息：
- 总标题：{outline.get('title', '')}
- 所有章节：
{chr(10).join([f"  - {outline.get(key, {}).get('title', '')}" for key in outline.keys() if key.startswith('section')])}
"""
                
                # 构建提示词
                prompt = f"""
                主题：{topic}
                
                这是PPT中的第 {slide_index + 1} 张幻灯片，以下是幻灯片的结构和内容信息：
                {json.dumps(slide_content_for_llm, ensure_ascii=False, indent=2)}
                
                {outline_info}
                
                请根据幻灯片的布局类型和原始内容，生成新的内容。要求：
                1. 保持相同的结构和格式
                2. 所有文本内容必须紧密围绕主题和当前章节主题
                3. 标题应该简短有力，不超过10个字
                4. 正文内容应该结构清晰，每段不超过150字
                5. 页脚、页码等信息保持原样
                6. 如果看到日期占位符，使用"{outline.get('date', '')}"替换
                7. 不允许返回原始内容信息，必须根据主题重新生成
                8. 如果只有感谢使用或者感谢观看或者Thanks这种结束语，就必须保持与原文本一致
                9. 请务必确保字数与原文本相近，不要相差太多
                10. 如果是表格，请确保表格的行数和列数与原始表格一致
                11. 内容必须符合当前章节的主题和内容要求
                """
                
                # 调用API生成内容
                response = self.client.chat.completions.create(
                    model=model_name,
                    messages=[
                        {"role": "system", "content": "你是一个专业的PPT内容生成助手。请确保生成的内容简洁精炼，适合PPT展示，并且整个页面的内容保持连贯性。内容必须严格遵循大纲的要求。"},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    response_format={"type": "json_object"}
                )
                
                # 解析生成的内容
                try:
                    generated_content = json.loads(response.choices[0].message.content)
                    logger.info(f"第 {slide_index + 1} 张幻灯片生成的内容：")
                    logger.info(json.dumps(generated_content, ensure_ascii=False, indent=2))
                    # 应用生成的内容到新幻灯片
                    self.apply_generated_content(new_prs.slides[slide_index], generated_content)
                except json.JSONDecodeError as e:
                    logger.error(f"JSON解析错误: {e}")
                    return {"error": f"内容生成格式错误: {e}"}
            
            # 保存生成的PPT
            new_prs.save(output_path)
            
            return {
                "success": True,
                "message": "PPT生成成功",
                "output_path": os.path.abspath(output_path)
            }
                
        except Exception as e:
            logger.error(f"生成PPT失败: {e}")
            return {"error": str(e)}

def main():
    """主函数：用于本地测试"""
    try:
        # 设置测试参数
        api_key = os.getenv("API_KEY")
        base_url = os.getenv("BASE_URL")
        print(api_key, base_url)
        template_path = "template_dfsj.pptx"
        output_path = "generated_ppt1.pptx"
        model_name = "Qwen-72B"
        topic = "民航局关于充电宝携带新规定"
        
        outline = {
            "title": "民航局关于充电宝携带新规定",
            "section1": {
                "title": "政策背景",
                "content": "介绍民航局发布新规的背景和原因，强调航空安全的重要性"
            },
            "section2": {
                "title": "了解最新航空安全政策",
                "content": "详细说明新规的具体要求：\n- 3C标识要求\n- 标识清晰度标准\n- 被召回型号限制"
            },
            "section3": {
                "title": "充电宝携带规则变更",
                "content": "明确新规实施日期：2023年6月28日"
            },
            "section4": {
                "title": "禁止携带无3C标识充电宝",
                "content": "说明新规适用范围：所有境内航班"
            },
            "section5": {
                "title": "实施时间与影响范围",
                "content": "为旅客提供实用建议：\n- 如何检查充电宝\n- 替代方案\n- 违规后果"
            },
            "date": "2024年3月21日"
        }
        
        # 初始化生成器
        generator = PPTGenerator(api_key=api_key, base_url=base_url)
        
        # 生成PPT
        result = generator.analyze_and_generate_ppt(
            template_path=template_path,
            topic=topic,
            output_path=output_path,
            model_name=model_name,
            outline=outline
        )
        
        if result.get("success"):
            print(f"PPT生成成功：{result['output_path']}")
        else:
            print(f"生成失败：{result.get('error')}")
            
    except Exception as e:
        print(f"程序执行失败: {e}")

if __name__ == '__main__':
    main() 