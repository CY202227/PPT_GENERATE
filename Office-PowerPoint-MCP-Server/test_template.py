#!/usr/bin/env python
"""
分析PPT模板并创建新的PPT文件
"""
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.shapes.group import GroupShape

def extract_shape_content(shape) -> dict:
    """递归提取形状内容"""
    element = None
    
    if isinstance(shape, GroupShape):
        elements = []
        for s in shape.shapes:
            child = extract_shape_content(s)
            if child:
                elements.append(child)
        if elements:
            element = {
                "type": "group",
                "elements": elements
            }
            
    elif hasattr(shape, 'text_frame') and shape.has_text_frame:
        text_frame = shape.text_frame
        element = {
            "type": "text",
            "content": text_frame.text
        }
        
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        element = {
            "type": "image",
            "image_info": {
                "filename": getattr(shape, "image_filename", ""),
                "size": getattr(shape, "image_size", "")
            }
        }
        
    elif hasattr(shape, 'table'):
        rows = []
        for row in shape.table.rows:
            cell_contents = [cell.text for cell in row.cells]
            rows.append(cell_contents)
        element = {
            "type": "table",
            "rows": rows
        }
    
    if element:
        element.update({
            "id": shape.shape_id,
            "name": shape.name,
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        })
    
    return element

def analyze_template(template_path: str) -> dict:
    """分析PPT模板并生成JSON结构"""
    prs = Presentation(template_path)
    template_data = {
        "slides": [],
        "layouts": []
    }
    
    # 分析幻灯片布局
    for layout in prs.slide_layouts:
        layout_data = {
            "name": layout.name,
            "elements": []
        }
        
        for shape in layout.shapes:
            element = extract_shape_content(shape)
            if element:
                layout_data["elements"].append(element)
                
        template_data["layouts"].append(layout_data)
    
    # 分析现有幻灯片
    for slide in prs.slides:
        slide_data = {
            "layout_name": slide.slide_layout.name,
            "elements": []
        }
        
        for shape in slide.shapes:
            element = extract_shape_content(shape)
            if element:
                slide_data["elements"].append(element)
                
        template_data["slides"].append(slide_data)
    
    return template_data

def create_new_presentation(template_data: dict, output_path: str):
    """使用MCP创建新的PPT文件"""
    from ppt_mcp_server import app, presentations
    
    # 初始化一个新的演示文稿
    prs = Presentation()
    presentation_id = "test_presentation"
    presentations[presentation_id] = prs
    
    # 添加幻灯片
    for slide_data in template_data["slides"]:
        # 找到对应的布局
        layout_name = slide_data["layout_name"]
        matching_layout = None
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                matching_layout = layout
                break
        
        # 如果找不到对应布局，使用空白布局
        if not matching_layout:
            matching_layout = prs.slide_layouts[6]  # 6通常是空白布局
        
        # 创建新幻灯片
        slide = prs.slides.add_slide(matching_layout)
        
        # 添加元素
        for element in slide_data["elements"]:
            if element["type"] == "text":
                shape = slide.shapes.add_textbox(
                    element["position"]["left"],
                    element["position"]["top"],
                    element["position"]["width"],
                    element["position"]["height"]
                )
                shape.text_frame.text = element["content"]
            
            elif element["type"] == "table" and "rows" in element:
                rows = element["rows"]
                if rows:
                    table = slide.shapes.add_table(
                        len(rows),
                        len(rows[0]),
                        element["position"]["left"],
                        element["position"]["top"],
                        element["position"]["width"],
                        element["position"]["height"]
                    ).table
                    
                    for i, row in enumerate(rows):
                        for j, cell_content in enumerate(row):
                            table.cell(i, j).text = cell_content
    
    # 保存PPT
    prs.save(output_path)

def main():
    # 设置路径
    template_path = "template_dfsj.pptx"
    output_json_path = "Office-PowerPoint-MCP-Server/slide_layout_templates.json"
    output_ppt_path = "generated_new.pptx"
    
    # 分析模板
    print("正在分析模板...")
    template_data = analyze_template(template_path)
    
    # 保存JSON
    print(f"正在保存模板数据到 {output_json_path}")
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(template_data, f, ensure_ascii=False, indent=2)
    
    # 创建新的PPT
    print(f"正在创建新的PPT文件 {output_ppt_path}")
    create_new_presentation(template_data, output_ppt_path)
    
    print("完成！")

if __name__ == '__main__':
    main()
