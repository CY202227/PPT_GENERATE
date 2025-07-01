#!/usr/bin/env python
"""
MCP Server for PowerPoint manipulation using python-pptx.
支持本地测试和MCP服务器两种模式。
"""
import os
import json
import logging
from typing import Dict, Any, Optional, List, Tuple
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.shapes.group import GroupShape
from pptx.dml.color import RGBColor
import utils.template_utils as template_utils
from openai import OpenAI
from tools.template_tools import register_template_tools

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 初始化FastMCP服务器
app = FastMCP(
    name="ppt-mcp-server",
    description="MCP Server for PowerPoint manipulation using python-pptx",
    version="1.0.0",
    log_level="INFO"
)

# 存储当前加载的演示文稿
presentations = {}
current_presentation_id = None

def get_current_presentation_id():
    return current_presentation_id

# 注册模板工具
register_template_tools(app, presentations, get_current_presentation_id)

def _extract_shape_content(shape: Any) -> Optional[Dict[str, Any]]:
    """
    递归提取形状内容，支持组合形状
    """
    element = None
    if isinstance(shape, GroupShape):
        elements = []
        for s in shape.shapes:
            child = _extract_shape_content(s)
            if child:
                elements.append(child)
        if not elements:
            return None
        element = { "type": "group", "elements": elements }

    elif hasattr(shape, 'text_frame') and shape.has_text_frame:
        text_frame = shape.text_frame
        content = text_frame.text
        element = { "type": "text", "content": content }

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        element = {"type": "image"}

    elif hasattr(shape, 'table'):
        rows = []
        for row in shape.table.rows:
            cell_contents = [cell.text for cell in row.cells]
            rows.append(cell_contents)
        element = {"type": "table", "rows": rows}
    
    if element:
        element["id"] = shape.shape_id
        # 添加一个可读的name，便于调试
        element["name"] = shape.name
        element["position"] = {"left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height}

    return element

def extract_slide_content_for_llm(slide) -> Dict[str, Any]:
    """
    将幻灯片转换为结构化的JSON格式给LLM，只包含ID、类型、位置和内容。
    """
    slide_content = {
        "layout_type": slide.slide_layout.name,
        "elements": []
    }
    for shape in slide.shapes:
        element_data = _extract_shape_content(shape)
        if element_data:
            slide_content["elements"].append(element_data)
    return slide_content

def find_shape_by_id(shapes, shape_id: int) -> Optional[Any]:
    """
    通过ID递归查找形状
    """
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if isinstance(shape, GroupShape):
            found = find_shape_by_id(shape.shapes, shape_id)
            if found:
                return found
    return None

def _get_new_content_map(elements: List[Dict[str, Any]]) -> Dict[int, str]:
    """
    递归遍历生成的elements，创建一个从shape_id到新内容的映射。
    """
    content_map = {}
    for element in elements:
        shape_id = element.get("id")
        if not shape_id:
            continue
        
        if element.get("type") == "text" and "content" in element:
            content_map[shape_id] = element["content"]
        elif element.get("type") == "table" and "rows" in element:
            content_map[shape_id] = element["rows"]
        
        if element.get("type") == "group" and "elements" in element:
            content_map.update(_get_new_content_map(element["elements"]))
            
    return content_map

def _apply_content_to_shapes(shapes: List[Any], content_map: Dict[int, Any]):
    """
    递归遍历幻灯片中的所有形状，并使用content_map更新文本。
    """
    for shape in shapes:
        if shape.shape_id in content_map:
            new_content = content_map[shape.shape_id]
            
            is_text_shape = hasattr(shape, 'has_text_frame') and shape.has_text_frame
            is_table_shape = hasattr(shape, 'has_table') and shape.has_table

            if is_text_shape and isinstance(new_content, str):
                text_frame = shape.text_frame
                if text_frame.paragraphs:
                    # 保留第一段的格式
                    p = text_frame.paragraphs[0]
                    
                    if p.runs:
                        # 更新第一个run的文本，保留其格式
                        p.runs[0].text = new_content
                        # 删除段落中多余的runs
                        for i in range(len(p.runs) - 1, 0, -1):
                            p._p.remove(p.runs[i]._r)
                    else:
                        # 如果段落没有run，则添加一个新的
                        p.add_run().text = new_content
                    
                    # 删除多余的段落
                    for i in range(len(text_frame.paragraphs) - 1, 0, -1):
                        p_to_remove = text_frame.paragraphs[i]
                        text_frame._txBody.remove(p_to_remove._p)
            
            elif is_table_shape and isinstance(new_content, list):
                table = shape.table
                for r, row_content in enumerate(new_content):
                    for c, cell_content in enumerate(row_content):
                        if r < len(table.rows) and c < len(table.columns):
                            cell = table.cell(r, c)
                            
                            print(f"更新表格内容{cell_content}")
                            if isinstance(cell_content, dict):
                                cell.text = str(cell_content.get("content", ""))
                            else:
                                cell.text = str(cell_content)

        if isinstance(shape, GroupShape):
            _apply_content_to_shapes(shape.shapes, content_map)

def apply_generated_content(slide, generated_content: Dict[str, Any], template_manager=None):
    """
    将生成的内容应用到幻灯片, 只更新content字段。
    """
    content_map = _get_new_content_map(generated_content.get("elements", []))
    _apply_content_to_shapes(slide.shapes, content_map)

def get_section_for_slide_index(slide_index: int, total_slides: int, outline: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    根据幻灯片索引获取对应的章节信息
    
    Args:
        slide_index: 当前幻灯片索引
        total_slides: 总幻灯片数
        outline: 大纲信息
        
    Returns:
        Optional[Dict[str, Any]]: 对应的章节信息
    """
    # 跳过封面（第1页）和目录（第2页）
    if slide_index < 2:
        return None
        
    # 获取所有章节
    section_keys = [key for key in outline.keys() if key.startswith('section')]
    if not section_keys:
        return None
        
    # 计算实际内容页数（去掉封面和目录）
    content_slides = total_slides - 2
    # 计算每个章节对应的页数
    slides_per_section = content_slides / len(section_keys)
    # 计算当前幻灯片对应的章节索引
    adjusted_index = slide_index - 2  # 调整索引，去掉前两页
    section_index = int(adjusted_index / slides_per_section)
    # 确保索引不超出范围
    section_index = min(section_index, len(section_keys) - 1)
    
    # 获取对应的章节信息
    section_key = section_keys[section_index]
    return outline.get(section_key)

def is_ending_slide(slide_content: Dict[str, Any]) -> Tuple[bool, Optional[str]]:
    """
    判断是否是结束页，并返回结束语
    
    Args:
        slide_content: 幻灯片内容
        
    Returns:
        Tuple[bool, Optional[str]]: (是否是结束页, 结束语内容)
    """
    ending_keywords = ["感谢", "谢谢", "thanks", "thank you", "the end"]
    
    for element in slide_content.get("elements", []):
        if element.get("type") == "text":
            content = element.get("content", "").lower().strip()
            # 如果文本很短（少于15个字符）且包含结束语关键词
            if len(content) < 15 and any(keyword in content.lower() for keyword in ending_keywords):
                return True, element.get("content")
            # 如果是单独的结束语（只包含结束语关键词和标点符号）
            if all(char in "！!。.，,、?？" or any(keyword in content.lower() for keyword in ending_keywords) for char in content):
                return True, element.get("content")
    return False, None

@app.tool()
def analyze_and_generate_ppt(
    template_path: str,
    topic: str,
    output_path: str,
    api_key: str,
    base_url: str,
    model_name: str,
    outline: Dict[str, Any] = None
) -> Dict[str, Any]:
    """
    分析PPT模板并生成新的PPT内容
    
    Args:
        template_path: PPT模板文件路径
        topic: 演示主题
        output_path: 输出PPT文件路径
        api_key: OpenAI API密钥
        base_url: OpenAI API基础URL
        model_name: 使用的模型名称
        outline: PPT大纲信息
        
    Returns:
        Dict[str, Any]: 包含生成结果的字典
    """
    try:
        if not os.path.exists(template_path):
            return {"error": f"模板文件不存在：{template_path}"}
            
        client = OpenAI(api_key=api_key, base_url=base_url)
        
        prs = Presentation(template_path)
        new_prs = Presentation(template_path)
        
        total_slides = len(prs.slides)
        
        # 遍历每个幻灯片
        for slide_index, slide in enumerate(prs.slides):
            logger.info(f"处理第 {slide_index + 1} 张幻灯片")
            
            # 提取幻灯片内容为JSON格式
            slide_content_for_llm = extract_slide_content_for_llm(slide)
            
            # 检查是否是最后一页，且是结束语
            is_ending, ending_text = is_ending_slide(slide_content_for_llm)
            if slide_index == total_slides - 1 and is_ending and ending_text:
                logger.info(f"检测到结束页，保持原有结束语：{ending_text}")
                # 直接复制原内容到新的幻灯片
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.has_text_frame:
                        shape.text_frame.text = shape.text_frame.text
                continue
            
            # 获取当前幻灯片对应的大纲部分
            current_section = None
            outline_info = ""
            if outline:
                current_section = get_section_for_slide_index(slide_index, total_slides, outline)
                if current_section:
                    outline_info = f"""
当前大纲信息：
- 总标题：{outline.get('title', '')}
- 当前章节：{current_section.get('title', '')}
- 章节内容：{current_section.get('content', '')}
- 日期：{outline.get('date', '')}
"""
                elif slide_index == 0:
                    outline_info = f"""
封面信息：
- 总标题：{outline.get('title', '')}
- 日期：{outline.get('date', '')}
"""
                elif slide_index == 1:
                    outline_info = f"""
目录页信息：
- 总标题：{outline.get('title', '')}
- 所有章节：
{chr(10).join([f"  - {outline.get(key, {}).get('title', '')}" for key in outline.keys() if key.startswith('section')])}
"""
            
            # 构建提示词
            prompt = f"""
            主题：{topic}
            
            这是PPT中的第 {slide_index + 1} 张幻灯片，以下是幻灯片的结构和内容信息：
            {json.dumps(slide_content_for_llm, ensure_ascii=False, indent=2)}
            
            {outline_info}
            
            请根据幻灯片的布局类型和原始内容，生成新的内容。要求：
            1. 保持相同的结构和格式
            2. 所有文本内容必须紧密围绕主题和当前章节主题
            3. 标题应该简短有力，不超过10个字
            4. 正文内容应该结构清晰，每段不超过150字
            5. 页脚、页码等信息保持原样
            6. 如果看到日期占位符，使用"{outline.get('date', '')}"替换
            7. 不允许返回原始内容信息，必须根据主题重新生成
            8. 如果只有 **感谢使用**或者**感谢观看**或者**Thanks**这种结束语，就必须保持与原文本一致,因为这是幻灯片的结束语
            9. 请务必确保字数与原文本相近，不要相差太多
            10. 如果是表格，请确保表格的行数和列数与原始表格一致
            11. 内容必须符合当前章节的主题和内容要求
            
            请返回一个JSON对象，保持与输入相同的结构，但将content字段替换为生成的新内容。
            其他字段（position、style等）保持不变。
            """
            
            # 调用API生成内容
            response = client.chat.completions.create(
                model=model_name,
                messages=[
                    {"role": "system", "content": "你是一个专业的PPT内容生成助手。请确保生成的内容简洁精炼，适合PPT展示，并且整个页面的内容保持连贯性。内容必须严格遵循大纲的要求。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                response_format={"type": "json_object"}
            )
            print(prompt)
            # 解析生成的内容
            try:
                generated_content = json.loads(response.choices[0].message.content)
                print(f"第 {slide_index + 1} 张幻灯片生成的内容：")
                print(json.dumps(generated_content, ensure_ascii=False, indent=2))
                # 应用生成的内容到新幻灯片
                apply_generated_content(new_prs.slides[slide_index], generated_content)
            except json.JSONDecodeError as e:
                logger.error(f"JSON解析错误: {e}")
                return {"error": f"内容生成格式错误: {e}"}
        
        # 保存生成的PPT
        new_prs.save(output_path)
        
        return {
            "success": True,
            "message": "PPT生成成功",
            "output_path": os.path.abspath(output_path)
        }
            
    except Exception as e:
        logger.error(f"生成PPT失败: {e}")
        return {"error": str(e)}

def main():
    """主函数：用于本地测试"""
    try:
        # 设置测试参数
        template_path = r"D:\Dev\PPT_Generate\template_dfsj.pptx"
        output_path = r"D:\Dev\PPT_Generate\generated_ppt1.pptx"
        topic = "日前，为切实保障航空运行安全，民航局发布紧急通知，自6月28日起禁止旅客携带没有3C标识、3C标识不清晰、被召回型号或批次的充电宝乘坐境内航班。"
        api_key = os.getenv("API_KEY")
        base_url = os.getenv("API_BASE")
        print(api_key, base_url)
        model_name = "Qwen-72B"
        outline = {
            "title": "民航局关于充电宝携带新规定",
            "section1": {
                "title": "政策背景",
                "content": "介绍民航局发布新规的背景和原因，强调航空安全的重要性"
            },
            "section2": {
                "title": "了解最新航空安全政策",
                "content": "详细说明新规的具体要求：\n- 3C标识要求\n- 标识清晰度标准\n- 被召回型号限制"
            },
            "section3": {
                "title": "充电宝携带规则变更",
                "content": "明确新规实施日期：2023年6月28日"
            },
            "section4": {
                "title": "禁止携带无3C标识充电宝",
                "content": "说明新规适用范围：所有境内航班"
            },
            "section5": {
                "title": "实施时间与影响范围",
                "content": "为旅客提供实用建议：\n- 如何检查充电宝\n- 替代方案\n- 违规后果"
            },
            "date": "2025年7月1日"
        }
        
        # 调用生成函数
        result = analyze_and_generate_ppt(
            template_path=template_path,
            topic=topic,
            output_path=output_path,
            api_key=api_key,
            base_url=base_url,
            model_name=model_name,
            outline=outline
        )
        
        if result.get("success"):
            print(f"\nPPT生成成功，输出文件: {result['output_path']}")
        else:
            print(f"\nPPT生成失败: {result.get('error')}")
    
    except Exception as e:
        print(f"\n程序执行失败: {e}")

if __name__ == '__main__':
    if os.getenv("MCP_MODE") == "server":
        # MCP服务器模式
        app.run()
    else:
        # 本地测试模式
        main()